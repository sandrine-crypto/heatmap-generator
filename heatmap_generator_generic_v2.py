#!/usr/bin/env python3
"""
Script générique ultra-flexible pour générer des heatmaps
Compatible avec n'importe quelle variable en première colonne :
- Temps (0h, 2h, D1, D3...)
- Concentrations (0, 10, 50, 100 ng/mL...)
- Doses, passages, cycles, etc.

Auteur: Assistant IA
Date: 2025-11-27
Version: 2.0 - Ultra-flexible
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import warnings
warnings.filterwarnings('ignore')


class HeatmapGenerator:
    """
    Générateur de heatmaps ultra-flexible pour données expérimentales
    """
    
    def __init__(self, fichier_csv, config=None):
        """
        Initialise le générateur
        
        Paramètres:
            fichier_csv (str): Chemin vers le fichier (CSV ou Excel)
            config (dict): Configuration optionnelle
        """
        self.fichier_csv = fichier_csv
        self.config = config or {}
        self.data = None
        self.marqueurs = []
        self.valeurs_x = []  # Anciennement temps_labels, maintenant générique
        self.groupes = []
        self.nom_colonne_x = None  # Nom détecté de la première colonne
        
        # Configuration par défaut
        self.default_config = {
            'colonne_x': None,  # None = détection automatique (première colonne)
            'label_axe_x': None,  # None = utilise le nom de la colonne
            'echelle_log': True,
            'taille_heatmap': (3.5, 4),
            'dpi': 150,
            'max_heatmaps_par_slide': 6,
            'titre_presentation': 'Analyse des résultats expérimentaux',
            'sous_titre': 'Heatmaps - Données quantitatives'
        }
        
        # Fusionner config
        for key, value in self.default_config.items():
            if key not in self.config:
                self.config[key] = value
    
    def charger_donnees(self):
        """
        Charge et analyse automatiquement la structure du fichier (CSV ou Excel)
        """
        print("\n" + "=" * 80)
        print("ANALYSE DU FICHIER DE DONNÉES")
        print("=" * 80 + "\n")
        
        # Déterminer le type de fichier et charger
        if self.fichier_csv.endswith('.xlsx') or self.fichier_csv.endswith('.xls'):
            print(f"Type de fichier détecté : Excel (.xlsx/.xls)")
            # Essayer d'abord de lire la feuille "Données", sinon la première feuille
            try:
                df = pd.read_excel(self.fichier_csv, sheet_name='Données', engine='openpyxl')
                print(f"  - Feuille utilisée : 'Données'")
            except:
                df = pd.read_excel(self.fichier_csv, sheet_name=0, engine='openpyxl')
                print(f"  - Feuille utilisée : première feuille")
        elif self.fichier_csv.endswith('.csv'):
            print(f"Type de fichier détecté : CSV")
            df = pd.read_csv(self.fichier_csv)
        else:
            raise ValueError(f"Format de fichier non supporté. Utilisez .xlsx, .xls ou .csv")
        
        print(f"✓ Fichier chargé : {self.fichier_csv}")
        print(f"  - Dimensions brutes : {df.shape[0]} lignes × {df.shape[1]} colonnes")
        
        # Détecter automatiquement la colonne X (première colonne) si non spécifiée
        if self.config['colonne_x'] is None:
            colonne_x = df.columns[0]
            print(f"  - Détection automatique : première colonne = '{colonne_x}'")
        else:
            colonne_x = self.config['colonne_x']
            print(f"  - Colonne X spécifiée : '{colonne_x}'")
        
        self.nom_colonne_x = colonne_x
        
        # Nettoyer les données
        if colonne_x in df.columns:
            # Supprimer les lignes où la colonne X est NaN ou vide
            df = df.dropna(subset=[colonne_x])
            
            # Supprimer les lignes avec texte long (instructions, notes)
            df = df[df[colonne_x].astype(str).str.len() < 20]
            
            # Supprimer les lignes avec emoji ou marqueurs d'instruction
            df = df[~df[colonne_x].astype(str).str.contains('📝|INSTRUCTION|Note|INFO|AIDE', case=False, na=False)]
        else:
            raise ValueError(f"Colonne '{colonne_x}' non trouvée. Colonnes disponibles : {list(df.columns)}")
        
        print(f"  - Dimensions nettoyées : {df.shape[0]} lignes × {df.shape[1]} colonnes")
        
        # Extraire les valeurs de la colonne X
        self.valeurs_x = df[colonne_x].tolist()
        print(f"  - {len(self.valeurs_x)} valeurs sur l'axe X : {self.valeurs_x}")
        
        # Identifier les colonnes de données (toutes sauf la première)
        colonnes_donnees = [col for col in df.columns if col != colonne_x]
        
        if not colonnes_donnees:
            raise ValueError("Aucune colonne de données trouvée (hormis la colonne X)")
        
        # Détecter marqueurs et groupes
        self.marqueurs = self._detecter_marqueurs(colonnes_donnees)
        self.groupes = self._detecter_groupes(colonnes_donnees, self.marqueurs)
        
        print(f"  - {len(self.marqueurs)} marqueurs détectés : {self.marqueurs}")
        print(f"  - {len(self.groupes)} groupes détectés : {self.groupes}")
        
        # Stocker les données
        self.data = df
    
    def _detecter_marqueurs(self, colonnes):
        """
        Détecte automatiquement les marqueurs depuis les noms de colonnes
        Format attendu : Marqueur_Groupe_Réplicat
        """
        marqueurs = set()
        
        for col in colonnes:
            # Essayer différents séparateurs
            for sep in ['_', '-', '.']:
                if sep in col:
                    # Le marqueur est le premier élément
                    marqueur = col.split(sep)[0]
                    marqueurs.add(marqueur)
                    break
        
        return sorted(list(marqueurs))
    
    def _detecter_groupes(self, colonnes, marqueurs):
        """
        Détecte automatiquement les groupes de traitement
        """
        groupes = set()
        
        for col in colonnes:
            # Essayer différents séparateurs
            for sep in ['_', '-', '.']:
                if sep in col:
                    parties = col.split(sep)
                    if len(parties) >= 2:
                        # Le groupe est l'élément entre le marqueur et le numéro de réplicat
                        marqueur = parties[0]
                        if marqueur in marqueurs:
                            groupe = parties[1]
                            # Vérifier que ce n'est pas juste un numéro (réplicat)
                            if not groupe.isdigit():
                                groupes.add(groupe)
                    break
        
        return sorted(list(groupes))
    
    def calculer_matrices(self):
        """
        Calcule les matrices de données pour chaque marqueur
        Format: [valeurs_x × groupes] avec les moyennes des réplicats
        """
        print("\n" + "=" * 80)
        print("CALCUL DES MATRICES DE DONNÉES")
        print("=" * 80 + "\n")
        
        self.matrices = {}
        colonne_x = self.nom_colonne_x
        
        for marqueur in self.marqueurs:
            # Initialiser matrice
            matrice = np.zeros((len(self.valeurs_x), len(self.groupes)))
            
            # Pour chaque groupe
            for j, groupe in enumerate(self.groupes):
                # Trouver toutes les colonnes pour ce marqueur et ce groupe
                colonnes_replicats = []
                for col in self.data.columns:
                    if col != colonne_x:
                        # Vérifier si la colonne correspond au marqueur et groupe
                        for sep in ['_', '-', '.']:
                            if sep in col:
                                parties = col.split(sep)
                                if len(parties) >= 2:
                                    if parties[0] == marqueur and parties[1] == groupe:
                                        colonnes_replicats.append(col)
                                break
                
                # Calculer moyennes pour chaque valeur X
                for i, val_x in enumerate(self.valeurs_x):
                    valeurs_replicats = []
                    for col_rep in colonnes_replicats:
                        val = self.data.loc[self.data[colonne_x] == val_x, col_rep].values
                        if len(val) > 0 and not np.isnan(val[0]):
                            valeurs_replicats.append(val[0])
                    
                    if valeurs_replicats:
                        matrice[i, j] = np.mean(valeurs_replicats)
                    else:
                        matrice[i, j] = 0
            
            self.matrices[marqueur] = matrice
            print(f"✓ {marqueur:<20} : matrice {len(self.valeurs_x)}×{len(self.groupes)}")
    
    def creer_heatmap(self, marqueur, afficher_valeurs=True, palette='rouge'):
        """
        Crée une heatmap pour un marqueur donné
        """
        matrice = self.matrices[marqueur]
        
        # Appliquer échelle logarithmique si demandé
        if self.config['echelle_log']:
            matrice_plot = np.log10(matrice + 1)
        else:
            matrice_plot = matrice
        
        # Définir les palettes de couleurs
        palettes = {
            'rouge': ['#FFFFFF', '#FFF5E6', '#FFE6CC', '#FFD9B3', '#FFCC99', 
                     '#FFB366', '#FF9933', '#FF8000', '#E67300', '#CC6600',
                     '#B35900', '#994C00', '#803F00'],
            'bleu': ['#FFFFFF', '#E6F2FF', '#CCE5FF', '#B3D9FF', '#99CCFF',
                    '#80BFFF', '#66B3FF', '#4DA6FF', '#3399FF', '#1A8CFF',
                    '#0080FF', '#0073E6', '#0066CC'],
            'vert': ['#FFFFFF', '#E6F9E6', '#CCF2CC', '#B3ECB3', '#99E699',
                    '#80DF80', '#66D966', '#4DD34D', '#33CC33', '#1AC61A',
                    '#00BF00', '#00B300', '#00A600'],
            'viridis': plt.cm.viridis(np.linspace(0, 1, 13)),
            'plasma': plt.cm.plasma(np.linspace(0, 1, 13))
        }
        
        # Créer colormap
        if palette in ['viridis', 'plasma']:
            cmap = LinearSegmentedColormap.from_list(palette, palettes[palette], N=100)
        else:
            cmap = LinearSegmentedColormap.from_list(palette, palettes[palette], N=100)
        
        # Créer figure
        fig, ax = plt.subplots(figsize=self.config['taille_heatmap'], dpi=self.config['dpi'])
        
        # Créer heatmap
        im = ax.imshow(matrice_plot, cmap=cmap, aspect='auto')
        
        # Configurer axes
        ax.set_xticks(np.arange(len(self.groupes)))
        ax.set_yticks(np.arange(len(self.valeurs_x)))
        ax.set_xticklabels(self.groupes, fontsize=9)
        ax.set_yticklabels(self.valeurs_x, fontsize=9)
        
        # Labels des axes
        label_x = self.config['label_axe_x'] if self.config['label_axe_x'] else self.nom_colonne_x
        ax.set_xlabel(label_x, fontsize=10, fontweight='bold')
        
        # Rotation des labels X si nécessaire
        plt.setp(ax.get_xticklabels(), rotation=45, ha="right", rotation_mode="anchor")
        
        # Afficher valeurs si demandé
        if afficher_valeurs:
            for i in range(len(self.valeurs_x)):
                for j in range(len(self.groupes)):
                    val = matrice[i, j]
                    # Format adaptatif
                    if val < 1:
                        text = f'{val:.2f}'
                    elif val < 10:
                        text = f'{val:.1f}'
                    else:
                        text = f'{val:.0f}'
                    
                    # Couleur du texte selon fond
                    if matrice_plot[i, j] > (matrice_plot.max() * 0.6):
                        color = 'white'
                    else:
                        color = 'black'
                    
                    ax.text(j, i, text, ha="center", va="center",
                           color=color, fontsize=8)
        
        # Titre
        ax.set_title(marqueur, fontsize=12, fontweight='bold', pad=10)
        
        # Colorbar
        cbar = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
        if self.config['echelle_log']:
            cbar.set_label('log₁₀(valeur + 1)', rotation=270, labelpad=15, fontsize=8)
        else:
            cbar.set_label('Valeur', rotation=270, labelpad=15, fontsize=8)
        
        plt.tight_layout()
        
        # Sauvegarder en mémoire
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', dpi=self.config['dpi'])
        buffer.seek(0)
        plt.close()
        
        return buffer
    
    def creer_presentation(self, fichier_sortie, afficher_valeurs=True, palette='rouge'):
        """
        Crée une présentation PowerPoint complète
        """
        print("\n" + "=" * 80)
        print("CRÉATION DE LA PRÉSENTATION")
        print("=" * 80 + "\n")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Page de titre
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = self.config['titre_presentation']
        p = tf.paragraphs[0]
        p.font.size = Inches(0.4)
        p.font.bold = True
        
        # Sous-titre
        top = Inches(3.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = self.config['sous_titre']
        p = tf.paragraphs[0]
        p.font.size = Inches(0.25)
        
        # Informations
        top = Inches(4.5)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        label_x = self.config['label_axe_x'] if self.config['label_axe_x'] else self.nom_colonne_x
        
        info_text = f"""• {len(self.marqueurs)} marqueurs analysés
• {len(self.valeurs_x)} valeurs de {label_x}
• {len(self.groupes)} groupes de traitement"""
        
        if self.config['echelle_log']:
            info_text += "\n• Échelle logarithmique appliquée"
        
        tf.text = info_text
        p = tf.paragraphs[0]
        p.font.size = Inches(0.18)
        
        print("✓ Slide 1 : Page de titre")
        
        # Slides suivantes: Heatmaps
        max_par_slide = self.config['max_heatmaps_par_slide']
        
        # Calculer disposition
        if max_par_slide <= 4:
            cols, rows = 2, 2
        elif max_par_slide <= 6:
            cols, rows = 3, 2
        else:
            cols, rows = 3, 3
        
        heatmap_width = 9 / cols
        heatmap_height = 6 / rows
        
        # Générer heatmaps
        num_slides_data = 0
        for i in range(0, len(self.marqueurs), max_par_slide):
            batch = self.marqueurs[i:i + max_par_slide]
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Titre du slide
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = f"Heatmaps - Marqueurs {i+1} à {min(i+max_par_slide, len(self.marqueurs))}"
            p = tf.paragraphs[0]
            p.font.size = Inches(0.25)
            p.font.bold = True
            
            # Ajouter heatmaps
            for j, marqueur in enumerate(batch):
                row = j // cols
                col = j % cols
                
                left = Inches(0.5 + col * heatmap_width)
                top = Inches(1 + row * heatmap_height)
                
                img_buffer = self.creer_heatmap(marqueur, afficher_valeurs, palette)
                slide.shapes.add_picture(img_buffer, left, top,
                                        width=Inches(heatmap_width * 0.9),
                                        height=Inches(heatmap_height * 0.85))
            
            num_slides_data += 1
        
        print(f"✓ Slide 2-{num_slides_data+1} : {len(self.marqueurs)} heatmaps ({', '.join(self.marqueurs)})")
        
        # Sauvegarder
        prs.save(fichier_sortie)
        
        print("\n" + "=" * 80)
        print(f"PRÉSENTATION CRÉÉE : {fichier_sortie}")
        print("=" * 80 + "\n")
        print(f"  - {num_slides_data + 1} slides générées")
        print(f"  - {len(self.marqueurs)} heatmaps au total")
        print(f"  - Palette : {palette}")
        print(f"  - Valeurs : {'affichées' if afficher_valeurs else 'masquées'}")
        print()


if __name__ == "__main__":
    # Exemple d'utilisation
    print("Ce module doit être importé, pas exécuté directement.")
    print("Utilisez: from heatmap_generator_generic import HeatmapGenerator")
